"""
Persona 3: Finance - Single combined storytelling slide.
Three-column flow: Platform View -> Department View -> Audit Trail
Bottom strip: The four numbers that answer the CFO's question.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import io

# Brand
DARK = RGBColor(0x15, 0x15, 0x15)
TEXT = RGBColor(0x43, 0x43, 0x43)
LIGHT = RGBColor(0x6A, 0x6E, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xEE, 0x00, 0x00)
TEAL = RGBColor(0x00, 0x97, 0xA7)
BLUE = RGBColor(0x00, 0x66, 0xCC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
CARD_BG = RGBColor(0xF8, 0xF8, 0xF8)
HEADER_BG = RGBColor(0x0E, 0x16, 0x1D)
H = "Red Hat Display"
B = "Red Hat Text"


def tb(s, l, t, w, h, text, sz=10, c=TEXT, bold=False, align=PP_ALIGN.LEFT, font=B):
    box = s.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Pt(0)
    box.text_frame.margin_right = Pt(0)
    box.text_frame.margin_top = Pt(0)
    box.text_frame.margin_bottom = Pt(0)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def card(s, l, t, w, h, fill=CARD_BG):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.5)
    return shape


def bar(s, l, t, w, color):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def circle_num(s, l, t, num, color):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.28), Inches(0.28))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    p = c.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = H
    p.alignment = PP_ALIGN.CENTER


def arrow(s, l, t):
    tb(s, l, t, Inches(0.25), Inches(0.25), "→", sz=16, c=LIGHT, bold=True, align=PP_ALIGN.CENTER)


# Load template for branding
template = Presentation("../rhoai-observability-external.pptx")
logo_blob = None
for shape in template.slides[1].shapes:
    if shape.shape_type == 13 and shape.left > Emu(9000000):
        logo_blob = shape.image.blob
        break

prs = Presentation("../rhoai-observability-external.pptx")
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

s = prs.slides.add_slide(prs.slide_layouts[6])

# ── Header bar ──
hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(12191695), Inches(0.5))
hdr.fill.solid()
hdr.fill.fore_color.rgb = HEADER_BG
hdr.line.fill.background()
tb(s, Inches(0.4), Inches(0.08), Inches(6), Inches(0.3),
   "FINANCE PERSONA", sz=8, c=RGBColor(0x99, 0x99, 0x99), bold=True, font=H)
tb(s, Inches(0.4), Inches(0.22), Inches(8), Inches(0.25),
   "\"What did AI cost us this quarter, by department?\"", sz=13, c=WHITE, bold=True, font=H)

# ── Three columns: Platform → Department → Audit ──
col_w = Inches(3.85)
col_gap = Inches(0.2)
col_y = Inches(0.7)
col_h = Inches(4.6)

columns = [
    {
        "num": 1,
        "color": RED,
        "title": "Platform view",
        "source": "Perses Usage Dashboard",
        "items": [
            ("837K tokens", "2K requests in 24 hours"),
            ("gemma4 primary", "538K tokens (64%)"),
            ("gemma4 secondary", "185K tokens (22%)"),
            ("qwen35-9b", "114K tokens (14%)"),
        ],
        "insight": "Three consumption streams. Engineering drove 64%. No setup required.",
    },
    {
        "num": 2,
        "color": TEAL,
        "title": "Department view",
        "source": "MLflow per-team experiments",
        "items": [
            ("Engineering", "305 traces, 61K tokens, 573 avg/trace"),
            ("Marketing", "169 traces, 28K tokens, 374 avg/trace"),
            ("Support", "109 traces, 15K tokens, 322 avg/trace"),
        ],
        "insight": "Engineering consumes 4x what Support does. That's the chargeback story.",
    },
    {
        "num": 3,
        "color": BLUE,
        "title": "Audit trail",
        "source": "MLflow Traces",
        "items": [
            ("Every request logged", "Trace ID, prompt, response, tokens"),
            ("Real prompts", "\"Review this code for vulnerabilities\""),
            ("Execution times", "5-6s avg, SLA baseline visible"),
        ],
        "insight": "Every inference call is traceable. The audit trail regulators require.",
    },
]

for i, col in enumerate(columns):
    x = Inches(0.35) + (col_w + col_gap) * i

    # Column card
    card(s, x, col_y, col_w, col_h)
    bar(s, x, col_y, col_w, col["color"])

    # Number + title
    circle_num(s, x + Inches(0.12), col_y + Inches(0.15), col["num"], col["color"])
    tb(s, x + Inches(0.48), col_y + Inches(0.15), Inches(2), Inches(0.25),
       col["title"], sz=14, c=DARK, bold=True, font=H)
    tb(s, x + Inches(0.48), col_y + Inches(0.4), Inches(3), Inches(0.2),
       col["source"], sz=8, c=LIGHT, font=B)

    # Data items
    item_y = col_y + Inches(0.75)
    for label, detail in col["items"]:
        # Label
        tb(s, x + Inches(0.15), item_y, col_w - Inches(0.3), Inches(0.2),
           label, sz=10, c=DARK, bold=True)
        # Detail
        tb(s, x + Inches(0.15), item_y + Inches(0.2), col_w - Inches(0.3), Inches(0.2),
           detail, sz=8, c=LIGHT)
        item_y += Inches(0.55)

    # Insight at bottom of column
    insight_y = col_y + col_h - Inches(0.65)
    ins_card = card(s, x + Inches(0.1), insight_y, col_w - Inches(0.2), Inches(0.5),
                    fill=WHITE)
    tb(s, x + Inches(0.2), insight_y + Inches(0.08), col_w - Inches(0.4), Inches(0.35),
       col["insight"], sz=8, c=col["color"], bold=False)

    # Arrows between columns
    if i < 2:
        arrow_x = x + col_w + Inches(0.02)
        arrow(s, arrow_x, col_y + Inches(2.0))

# ── Bottom strip: The Answer ──
strip_y = Inches(5.5)
strip_h = Inches(1.25)

strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, strip_y, Emu(12191695), strip_h)
strip.fill.solid()
strip.fill.fore_color.rgb = HEADER_BG
strip.line.fill.background()

tb(s, Inches(0.4), strip_y + Inches(0.08), Inches(3), Inches(0.2),
   "THE CFO'S ANSWER", sz=8, c=RGBColor(0x99, 0x99, 0x99), bold=True, font=H)

stats = [
    ("837K", "total tokens", RED),
    ("64%", "Engineering", TEAL),
    ("22%", "Marketing", BLUE),
    ("14%", "Support", RGBColor(0x63, 0x99, 0x3D)),
]

for i, (num, label, color) in enumerate(stats):
    x = Inches(0.5) + Inches(3.1) * i
    tb(s, x, strip_y + Inches(0.3), Inches(2), Inches(0.5),
       num, sz=32, c=color, bold=True, font=H)
    tb(s, x, strip_y + Inches(0.85), Inches(2.5), Inches(0.2),
       label, sz=10, c=RGBColor(0xBB, 0xBB, 0xBB))

# ── Footer ──
tb(s, Inches(0.4), Inches(6.95), Inches(8), Inches(0.2),
   "Red Hat OpenShift AI  |  Observability  |  Finance Persona Demo", sz=7, c=LIGHT)
if logo_blob:
    s.shapes.add_picture(io.BytesIO(logo_blob), Inches(11.2), Inches(6.9), Inches(1.0), Inches(0.23))

# Speaker notes
notes = s.notes_slide
notes.notes_text_frame.text = (
    "This slide tells the complete chargeback story in three steps.\n\n"
    "Step 1 (Platform view): Open the Perses Usage dashboard. 837K tokens, 2K requests. "
    "Three rows in the Token Consumption table showing consumption by subscription and model. "
    "Engineering drove 64% of tokens on the expensive model. Support used 14% on the cheap model.\n\n"
    "Step 2 (Department view): Switch to MLflow. Three experiments, one per department. "
    "Engineering consumed 305 traces at 573 tokens avg. Support consumed 109 traces at 322 avg. "
    "Engineering uses 4x what Support does. That's the cost attribution data.\n\n"
    "Step 3 (Audit trail): Click into the Engineering traces. 305 individual requests with full "
    "prompt and response content, token counts, and execution times. Every request traceable. "
    "This is the audit trail telco and FSI regulators require.\n\n"
    "Bottom strip: the four numbers that answer the CFO's question. 837K total, "
    "64% engineering, 22% marketing, 14% support. Today this is in dashboards. "
    "In the next release, the metering webhook pushes it to your billing system automatically."
)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: Persona 1 — Platform Admin / SRE
# "The 2am Latency Spike"
# ══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
hdr2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(12191695), Inches(0.5))
hdr2.fill.solid()
hdr2.fill.fore_color.rgb = HEADER_BG
hdr2.line.fill.background()
tb(s2, Inches(0.4), Inches(0.08), Inches(6), Inches(0.3),
   "PLATFORM ADMIN / SRE PERSONA", sz=8, c=RGBColor(0x99, 0x99, 0x99), bold=True, font=H)
tb(s2, Inches(0.4), Inches(0.22), Inches(10), Inches(0.25),
   "\"It's 2am. Latency spiked. What happened and how do I fix it?\"",
   sz=13, c=WHITE, bold=True, font=H)

# Five-step horizontal flow
step_w = Inches(2.25)
step_h = Inches(4.3)
step_gap = Inches(0.15)
step_y = Inches(0.7)

steps = [
    {
        "num": 1,
        "color": RGBColor(0x63, 0x99, 0x3D),
        "title": "Normal ops",
        "source": "Fed Aura chatbot",
        "items": [
            ("Chatbot responding", "< 1 second per request"),
            ("TTFT stable", "p99 at 200ms"),
            ("GPU at 45%", "Healthy headroom"),
        ],
        "insight": "Everything looks good. Baseline established.",
    },
    {
        "num": 2,
        "color": RGBColor(0xF0, 0x56, 0x1D),
        "title": "Load ramps",
        "source": "Traffic increases 3x",
        "items": [
            ("Concurrent users", "spike from 5 to 50"),
            ("Queue depth", "waiting > running"),
            ("TTFT climbing", "200ms → 2s → 4s"),
        ],
        "insight": "Something is wrong. Alert fires: p99 TTFT > 2s.",
    },
    {
        "num": 3,
        "color": RED,
        "title": "Root cause",
        "source": "Perses diagnostic panels",
        "items": [
            ("KV Cache at 100%", "Cache panel shows saturation"),
            ("GPU memory full", "No room for new KV entries"),
            ("Preemptions rising", "Scheduler thrashing"),
        ],
        "insight": "Cache exhaustion. Requests evicting each other.",
    },
    {
        "num": 4,
        "color": BLUE,
        "title": "Correlate",
        "source": "Korrel8r + Loki",
        "items": [
            ("OOM warnings", "From the same pod + time"),
            ("Logs correlated", "One query, not 4 tools"),
            ("Pod 2 identified", "Specific pod isolated"),
        ],
        "insight": "Korrel8r linked the metric alert to the OOM logs.",
    },
    {
        "num": 5,
        "color": TEAL,
        "title": "Fix",
        "source": "Action taken",
        "items": [
            ("Scale replicas", "Add pod to spread load"),
            ("Or reduce cache", "Lower max_model_len"),
            ("TTFT recovers", "Back to 200ms in minutes"),
        ],
        "insight": "2 minutes from alert to diagnosis. Not 45.",
    },
]

for i, step in enumerate(steps):
    x = Inches(0.25) + (step_w + step_gap) * i

    # Step card
    card(s2, x, step_y, step_w, step_h)
    bar(s2, x, step_y, step_w, step["color"])

    # Number + title
    circle_num(s2, x + Inches(0.1), step_y + Inches(0.12), step["num"], step["color"])
    tb(s2, x + Inches(0.42), step_y + Inches(0.12), Inches(1.5), Inches(0.22),
       step["title"], sz=13, c=DARK, bold=True, font=H)
    tb(s2, x + Inches(0.42), step_y + Inches(0.35), Inches(1.7), Inches(0.18),
       step["source"], sz=7, c=LIGHT)

    # Data items
    item_y = step_y + Inches(0.65)
    for label, detail in step["items"]:
        tb(s2, x + Inches(0.1), item_y, step_w - Inches(0.2), Inches(0.18),
           label, sz=9, c=DARK, bold=True)
        tb(s2, x + Inches(0.1), item_y + Inches(0.18), step_w - Inches(0.2), Inches(0.18),
           detail, sz=7, c=LIGHT)
        item_y += Inches(0.5)

    # Insight
    insight_y = step_y + step_h - Inches(0.6)
    card(s2, x + Inches(0.08), insight_y, step_w - Inches(0.16), Inches(0.45), fill=WHITE)
    tb(s2, x + Inches(0.15), insight_y + Inches(0.06), step_w - Inches(0.3), Inches(0.32),
       step["insight"], sz=7, c=step["color"])

    # Arrow between steps
    if i < 4:
        arrow_x = x + step_w + Inches(0.0)
        arrow(s2, arrow_x, step_y + Inches(1.8))

# Bottom strip
strip2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.2), Emu(12191695), Inches(1.55))
strip2.fill.solid()
strip2.fill.fore_color.rgb = HEADER_BG
strip2.line.fill.background()

tb(s2, Inches(0.4), Inches(5.28), Inches(4), Inches(0.18),
   "THE SRE'S TOOLKIT", sz=8, c=RGBColor(0x99, 0x99, 0x99), bold=True, font=H)

tools = [
    ("54+", "vLLM metrics\nout of the box", RED),
    ("6", "pre-built SLO\nalert rules", TEAL),
    ("10", "Perses diagnostic\npanels", BLUE),
    ("1 query", "Korrel8r correlates\nmetrics + logs + traces", RGBColor(0x63, 0x99, 0x3D)),
]

for i, (num, label, color) in enumerate(tools):
    x = Inches(0.5) + Inches(3.1) * i
    tb(s2, x, Inches(5.5), Inches(2), Inches(0.45),
       num, sz=28, c=color, bold=True, font=H)
    tb(s2, x, Inches(5.95), Inches(2.5), Inches(0.4),
       label, sz=9, c=RGBColor(0xBB, 0xBB, 0xBB))

# Footer
tb(s2, Inches(0.4), Inches(6.95), Inches(8), Inches(0.2),
   "Red Hat OpenShift AI  |  Observability  |  SRE Persona Demo", sz=7, c=LIGHT)
if logo_blob:
    s2.shapes.add_picture(io.BytesIO(logo_blob), Inches(11.2), Inches(6.9), Inches(1.0), Inches(0.23))

# Speaker notes
notes2 = s2.notes_slide
notes2.notes_text_frame.text = (
    "This slide walks through the 2am latency spike scenario in five steps.\n\n"
    "Step 1: everything is normal. The Fed Aura chatbot responds in under a second. "
    "TTFT is stable at 200ms. GPU utilization is at 45%. This is the baseline.\n\n"
    "Step 2: traffic ramps. Concurrent users go from 5 to 50. The queue depth panel "
    "shows waiting requests exceeding running requests. TTFT climbs from 200ms to 4s. "
    "The PrometheusRule fires: p99 TTFT exceeded 2 seconds for 5 minutes.\n\n"
    "Step 3: the SRE opens Perses. The KV cache panel shows 100% utilization. "
    "GPU memory is full. Preemptions are rising because the scheduler is thrashing, "
    "evicting cache entries to make room. This is the root cause.\n\n"
    "Step 4: Korrel8r correlates. One query finds OOM warning logs in Loki from the "
    "same pod and time window. Pod 2 is the problem. The SRE didn't have to open "
    "Prometheus, then Loki, then kubectl. One query across all backends.\n\n"
    "Step 5: the fix. Scale replicas to spread load, or reduce max_model_len to lower "
    "cache pressure. TTFT recovers in minutes.\n\n"
    "Bottom line: without observability, this is a 45 minute investigation across 4 tools. "
    "With Perses diagnostic panels and Korrel8r cross-signal correlation, the SRE identified "
    "root cause in under 2 minutes. That's the value proposition for the platform team."
)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: Persona 4 — CIO / Decision Maker
# "Ask Your AI Platform a Business Question"
# ══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
hdr3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(12191695), Inches(0.5))
hdr3.fill.solid()
hdr3.fill.fore_color.rgb = HEADER_BG
hdr3.line.fill.background()
tb(s3, Inches(0.4), Inches(0.08), Inches(6), Inches(0.3),
   "CIO / DECISION MAKER PERSONA", sz=8, c=RGBColor(0x99, 0x99, 0x99), bold=True, font=H)
tb(s3, Inches(0.4), Inches(0.22), Inches(10), Inches(0.25),
   "\"What's our ROI on GPU investment? Just ask the platform.\"",
   sz=13, c=WHITE, bold=True, font=H)

# Layout: Left side = MCP Server concept (the "brain")
# Right side = 5 business questions with answers

# Left panel: MCP Server
left_w = Inches(3.8)
left_x = Inches(0.3)
panel_y = Inches(0.7)
panel_h = Inches(4.5)

card(s3, left_x, panel_y, left_w, panel_h, fill=RGBColor(0xF0, 0xF0, 0xF0))
bar(s3, left_x, panel_y, left_w, RED)

# MCP Server title
tb(s3, left_x + Inches(0.2), panel_y + Inches(0.2), Inches(3.3), Inches(0.25),
   "MCP Server", sz=18, c=DARK, bold=True, font=H)
tb(s3, left_x + Inches(0.2), panel_y + Inches(0.5), Inches(3.3), Inches(0.2),
   "18 tools  |  5 backends  |  natural language", sz=9, c=LIGHT)

# The 5 backends as small pills
backends = [
    ("Prometheus", RED),
    ("AlertManager", RGBColor(0xF0, 0x56, 0x1D)),
    ("Loki", TEAL),
    ("Kubernetes", BLUE),
    ("Tempo", RGBColor(0x5E, 0x40, 0xBE)),
]
by = panel_y + Inches(0.9)
for name, color in backends:
    pill_shape = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left_x + Inches(0.2), by, Inches(1.5), Inches(0.25))
    pill_shape.fill.solid()
    pill_shape.fill.fore_color.rgb = color
    pill_shape.line.fill.background()
    p = pill_shape.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(8)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = B
    p.alignment = PP_ALIGN.CENTER
    by += Inches(0.32)

# Concept text
tb(s3, left_x + Inches(0.2), panel_y + Inches(2.7), Inches(3.3), Inches(0.6),
   "The CIO doesn't want a dashboard.\nThey want answers.", sz=11, c=DARK, bold=True, font=H)
tb(s3, left_x + Inches(0.2), panel_y + Inches(3.4), Inches(3.3), Inches(0.8),
   "Connect an AI assistant to your cluster's observability stack. "
   "Ask questions in plain language. Get answers with data, not PromQL.",
   sz=9, c=LIGHT)

# Right side: Q&A flow
right_x = Inches(4.35)
right_w = Inches(8.5)
qa_y = panel_y

questions = [
    {
        "q": "What's our GPU utilization across the cluster?",
        "a": "72% average. Node-3 at 95%. Node-7 at 31%.",
        "color": RED,
    },
    {
        "q": "Which models are getting the most traffic?",
        "a": "gemma4 handles 60% of requests. qwen35-9b handles 25%. llama takes 15%.",
        "color": TEAL,
    },
    {
        "q": "What would happen if we removed one GPU node?",
        "a": "Node-7 is underloaded. Removing it pushes avg to 85%, still within SLO.",
        "color": BLUE,
    },
    {
        "q": "Are we meeting our latency commitments?",
        "a": "p99 TTFT is 1.2s against a 2s SLA. Margin is healthy.",
        "color": RGBColor(0x63, 0x99, 0x3D),
    },
    {
        "q": "Generate a summary for the board.",
        "a": "Formatted report: utilization, traffic, SLA compliance, and recommendations.",
        "color": RGBColor(0x5E, 0x40, 0xBE),
    },
]

for i, qa in enumerate(questions):
    y = qa_y + Inches(0.9) * i
    qa_card_w = Inches(8.35)

    # Question/Answer row
    card(s3, right_x, y, qa_card_w, Inches(0.78))
    bar(s3, right_x, y, Pt(3), qa["color"])

    # Left accent bar (vertical)
    vbar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y, Pt(3), Inches(0.78))
    vbar.fill.solid()
    vbar.fill.fore_color.rgb = qa["color"]
    vbar.line.fill.background()

    # Q icon
    q_circle = s3.shapes.add_shape(MSO_SHAPE.OVAL,
                                    right_x + Inches(0.12), y + Inches(0.08),
                                    Inches(0.2), Inches(0.2))
    q_circle.fill.solid()
    q_circle.fill.fore_color.rgb = qa["color"]
    q_circle.line.fill.background()
    qp = q_circle.text_frame.paragraphs[0]
    qp.text = "Q"
    qp.font.size = Pt(9)
    qp.font.color.rgb = WHITE
    qp.font.bold = True
    qp.font.name = H
    qp.alignment = PP_ALIGN.CENTER

    # Question text
    tb(s3, right_x + Inches(0.4), y + Inches(0.07), Inches(7.5), Inches(0.2),
       qa["q"], sz=10, c=DARK, bold=True)

    # Answer text
    tb(s3, right_x + Inches(0.4), y + Inches(0.35), Inches(7.5), Inches(0.35),
       qa["a"], sz=9, c=LIGHT)

# Bottom strip
strip3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.5), Emu(12191695), Inches(1.25))
strip3.fill.solid()
strip3.fill.fore_color.rgb = HEADER_BG
strip3.line.fill.background()

tb(s3, Inches(0.4), Inches(5.58), Inches(4), Inches(0.18),
   "WHY THIS MATTERS", sz=8, c=RGBColor(0x99, 0x99, 0x99), bold=True, font=H)

points = [
    ("No PromQL", "Business questions,\nnot query languages", RED),
    ("5 backends", "One conversation,\nnot five tools", TEAL),
    ("Board-ready", "Export summaries\nfor stakeholders", BLUE),
    ("First mover", "No competitor\noffers this today", RGBColor(0x63, 0x99, 0x3D)),
]

for i, (num, label, color) in enumerate(points):
    x = Inches(0.5) + Inches(3.1) * i
    tb(s3, x, Inches(5.8), Inches(2), Inches(0.35),
       num, sz=20, c=color, bold=True, font=H)
    tb(s3, x, Inches(6.15), Inches(2.5), Inches(0.4),
       label, sz=9, c=RGBColor(0xBB, 0xBB, 0xBB))

# Footer
tb(s3, Inches(0.4), Inches(6.95), Inches(8), Inches(0.2),
   "Red Hat OpenShift AI  |  Observability  |  CIO Persona Demo", sz=7, c=LIGHT)
if logo_blob:
    s3.shapes.add_picture(io.BytesIO(logo_blob), Inches(11.2), Inches(6.9), Inches(1.0), Inches(0.23))

# Speaker notes
notes3 = s3.notes_slide
notes3.notes_text_frame.text = (
    "This slide shows the ACT pillar in action for the CIO persona.\n\n"
    "Left side: the MCP Server sits on top of five observability backends. "
    "Prometheus for metrics, AlertManager for alerts, Loki for logs, Kubernetes "
    "for cluster state, Tempo for traces. 18 tools, all queryable through "
    "natural language.\n\n"
    "Right side: five questions a CIO actually asks, with the answers the MCP "
    "Server provides. These aren't hypothetical. The AI Observability Summarizer "
    "running on the cluster today can answer all five.\n\n"
    "GPU utilization: tells the CIO if hardware investment is being used. "
    "72% average is good. Node-7 at 31% is waste they can cut.\n\n"
    "Model traffic: tells them which models justify their allocation. "
    "If gemma4 handles 60% of requests, that's where investment should go.\n\n"
    "Node removal: capacity planning. Can we save money by removing a node "
    "without breaking SLAs? The answer is data-driven, not a guess.\n\n"
    "SLA compliance: are we meeting commitments? p99 at 1.2s against a 2s target "
    "means healthy margin.\n\n"
    "Board summary: the CIO doesn't want to screenshot dashboards for their board deck. "
    "They want a formatted summary they can share. The MCP Server generates it.\n\n"
    "Bottom line: no competitor offers this. The CIO doesn't want a dashboard. "
    "They want answers. This is our first-mover advantage."
)

out = "demo-chargeback-story.pptx"
prs.save(out)
print(f"Done. 3 slides saved to {out}")
